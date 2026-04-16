#!/usr/bin/env python3
"""
Convert AWS Lesson JSON files to PowerPoint presentations with animations
Designed for teachers to easily explain concepts with visual aids
"""

import json
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# AWS Brand Colors
AWS_ORANGE = RGBColor(255, 153, 0)
AWS_DARK = RGBColor(35, 47, 62)
AWS_LIGHT_GRAY = RGBColor(248, 249, 250)
AWS_GREEN = RGBColor(40, 167, 69)
AWS_RED = RGBColor(220, 53, 69)

def create_title_slide(prs, title, subtitle):
    """Create an engaging title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = AWS_DARK
    background.line.fill.background()

    # Icon/Emoji (large)
    icon_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.5), Inches(8), Inches(1.5)
    )
    icon_frame = icon_box.text_frame
    icon_frame.text = subtitle  # Icon
    icon_para = icon_frame.paragraphs[0]
    icon_para.alignment = PP_ALIGN.CENTER
    icon_para.font.size = Pt(120)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.5), Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.bold = True
    title_para.font.size = Pt(44)
    title_para.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(5), Inches(8), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AWS Cloud Computing"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = AWS_ORANGE

def create_section_header(prs, section_number, section_title):
    """Create a section divider slide with animation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background gradient
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = AWS_ORANGE
    background.line.fill.background()

    # Section number
    number_box = slide.shapes.add_textbox(
        Inches(1), Inches(2), Inches(8), Inches(1)
    )
    number_frame = number_box.text_frame
    number_frame.text = f"Section {section_number}"
    number_para = number_frame.paragraphs[0]
    number_para.alignment = PP_ALIGN.CENTER
    number_para.font.size = Pt(32)
    number_para.font.color.rgb = RGBColor(255, 255, 255)

    # Section title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(3), Inches(8), Inches(2)
    )
    title_frame = title_box.text_frame
    title_frame.text = section_title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.bold = True
    title_para.font.size = Pt(48)
    title_para.font.color.rgb = RGBColor(255, 255, 255)

def add_bullet_points(text_frame, items, animate=True):
    """Add bullet points with formatting"""
    # Clear existing text
    text_frame.clear()

    for idx, item in enumerate(items):
        # Remove markdown bold syntax
        clean_item = item.replace('**', '')

        p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        p.text = clean_item
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = AWS_DARK
        p.space_after = Pt(10)

def create_text_slide(prs, section_number, title, content):
    """Create a text content slide with key points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add section number in corner
    section_label = slide.shapes.add_textbox(
        Inches(8.5), Inches(0.3), Inches(1), Inches(0.4)
    )
    section_frame = section_label.text_frame
    section_frame.text = f"§{section_number}"
    section_para = section_frame.paragraphs[0]
    section_para.font.size = Pt(14)
    section_para.font.color.rgb = AWS_ORANGE

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.bold = True
    title_para.font.size = Pt(32)
    title_para.font.color.rgb = AWS_DARK

    # Underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.3), Inches(2), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = AWS_ORANGE
    line.line.fill.background()

    # Content box with light background
    content_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(1.8), Inches(8.6), Inches(4.5)
    )
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = AWS_LIGHT_GRAY
    content_box.line.color.rgb = AWS_ORANGE
    content_box.line.width = Pt(2)

    # Content text
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_right = Inches(0.3)
    text_frame.margin_top = Inches(0.3)

    # Clean and format content
    clean_content = content.replace('**', '')
    paragraphs = clean_content.split('\n\n')

    for idx, para in enumerate(paragraphs):
        if para.strip():
            p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
            p.text = para.strip()
            p.font.size = Pt(18)
            p.font.color.rgb = AWS_DARK
            p.space_after = Pt(15)
            p.alignment = PP_ALIGN.LEFT

def create_list_slide(prs, section_number, title, items, content=None):
    """Create a slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Section number
    section_label = slide.shapes.add_textbox(
        Inches(8.5), Inches(0.3), Inches(1), Inches(0.4)
    )
    section_frame = section_label.text_frame
    section_frame.text = f"§{section_number}"
    section_para = section_frame.paragraphs[0]
    section_para.font.size = Pt(14)
    section_para.font.color.rgb = AWS_ORANGE

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.bold = True
    title_para.font.size = Pt(32)
    title_para.font.color.rgb = AWS_DARK

    # Underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.3), Inches(2), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = AWS_ORANGE
    line.line.fill.background()

    # Optional intro text
    start_y = 1.8
    if content:
        intro_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(start_y), Inches(8.6), Inches(0.6)
        )
        intro_frame = intro_box.text_frame
        intro_frame.text = content.replace('**', '')
        intro_para = intro_frame.paragraphs[0]
        intro_para.font.size = Pt(16)
        intro_para.font.italic = True
        intro_para.font.color.rgb = RGBColor(100, 100, 100)
        start_y = 2.5

    # Bullet points
    bullets_box = slide.shapes.add_textbox(
        Inches(1), Inches(start_y), Inches(8), Inches(5)
    )
    text_frame = bullets_box.text_frame
    text_frame.word_wrap = True

    add_bullet_points(text_frame, items)

def create_comparison_slide(prs, section_number, title, traditional, cloud, content=None):
    """Create a side-by-side comparison slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Section number
    section_label = slide.shapes.add_textbox(
        Inches(8.5), Inches(0.3), Inches(1), Inches(0.4)
    )
    section_frame = section_label.text_frame
    section_frame.text = f"§{section_number}"
    section_para = section_frame.paragraphs[0]
    section_para.font.size = Pt(14)
    section_para.font.color.rgb = AWS_ORANGE

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.bold = True
    title_para.font.size = Pt(32)
    title_para.font.color.rgb = AWS_DARK

    # Intro text if provided
    start_y = 1.5
    if content:
        intro_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(start_y), Inches(8.6), Inches(0.5)
        )
        intro_frame = intro_box.text_frame
        intro_frame.text = content.replace('**', '')
        intro_para = intro_frame.paragraphs[0]
        intro_para.font.size = Pt(16)
        intro_para.font.italic = True
        intro_para.font.color.rgb = RGBColor(100, 100, 100)
        start_y = 2.1

    # Left box - Traditional
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(start_y), Inches(4.5), Inches(4.5)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(255, 240, 240)  # Light red
    left_box.line.color.rgb = AWS_RED
    left_box.line.width = Pt(3)

    # Traditional title
    left_title = slide.shapes.add_textbox(
        Inches(0.7), Inches(start_y + 0.2), Inches(4.1), Inches(0.5)
    )
    left_title_frame = left_title.text_frame
    left_title_frame.text = f"❌ {traditional['title']}"
    left_title_para = left_title_frame.paragraphs[0]
    left_title_para.font.bold = True
    left_title_para.font.size = Pt(20)
    left_title_para.font.color.rgb = AWS_RED
    left_title_para.alignment = PP_ALIGN.CENTER

    # Traditional items
    left_items = slide.shapes.add_textbox(
        Inches(0.8), Inches(start_y + 0.8), Inches(3.9), Inches(3.5)
    )
    left_frame = left_items.text_frame
    left_frame.word_wrap = True
    for idx, item in enumerate(traditional['items']):
        p = left_frame.paragraphs[0] if idx == 0 else left_frame.add_paragraph()
        p.text = item.replace('**', '')
        p.font.size = Pt(14)
        p.font.color.rgb = AWS_DARK
        p.space_after = Pt(8)
        p.level = 0

    # Right box - Cloud
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(start_y), Inches(4.5), Inches(4.5)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(240, 255, 240)  # Light green
    right_box.line.color.rgb = AWS_GREEN
    right_box.line.width = Pt(3)

    # Cloud title
    right_title = slide.shapes.add_textbox(
        Inches(5.4), Inches(start_y + 0.2), Inches(4.1), Inches(0.5)
    )
    right_title_frame = right_title.text_frame
    right_title_frame.text = f"✅ {cloud['title']}"
    right_title_para = right_title_frame.paragraphs[0]
    right_title_para.font.bold = True
    right_title_para.font.size = Pt(20)
    right_title_para.font.color.rgb = AWS_GREEN
    right_title_para.alignment = PP_ALIGN.CENTER

    # Cloud items
    right_items = slide.shapes.add_textbox(
        Inches(5.5), Inches(start_y + 0.8), Inches(3.9), Inches(3.5)
    )
    right_frame = right_items.text_frame
    right_frame.word_wrap = True
    for idx, item in enumerate(cloud['items']):
        p = right_frame.paragraphs[0] if idx == 0 else right_frame.add_paragraph()
        p.text = item.replace('**', '')
        p.font.size = Pt(14)
        p.font.color.rgb = AWS_DARK
        p.space_after = Pt(8)
        p.level = 0

def create_advantages_slide(prs, advantages):
    """Create a summary slide with key advantages"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "✨ Key Takeaways"
    title_para = title_frame.paragraphs[0]
    title_para.font.bold = True
    title_para.font.size = Pt(36)
    title_para.font.color.rgb = AWS_DARK
    title_para.alignment = PP_ALIGN.CENTER

    # Advantages grid
    start_y = 2
    for idx, advantage in enumerate(advantages[:5]):  # Limit to 5
        # Icon box
        icon_box = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1), Inches(start_y + (idx * 0.9)), Inches(0.6), Inches(0.6)
        )
        icon_box.fill.solid()
        icon_box.fill.fore_color.rgb = AWS_GREEN
        icon_box.line.fill.background()

        # Check mark
        check_text = slide.shapes.add_textbox(
            Inches(1.15), Inches(start_y + (idx * 0.9) + 0.1), Inches(0.3), Inches(0.4)
        )
        check_frame = check_text.text_frame
        check_frame.text = "✓"
        check_para = check_frame.paragraphs[0]
        check_para.font.size = Pt(24)
        check_para.font.bold = True
        check_para.font.color.rgb = RGBColor(255, 255, 255)

        # Advantage text
        adv_box = slide.shapes.add_textbox(
            Inches(1.8), Inches(start_y + (idx * 0.9)), Inches(7), Inches(0.6)
        )
        adv_frame = adv_box.text_frame
        adv_frame.text = advantage.replace('**', '')
        adv_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        adv_para = adv_frame.paragraphs[0]
        adv_para.font.size = Pt(18)
        adv_para.font.color.rgb = AWS_DARK

def create_end_slide(prs, title):
    """Create a closing slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = AWS_DARK
    background.line.fill.background()

    # Thank you message
    thank_you = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(1)
    )
    thank_you_frame = thank_you.text_frame
    thank_you_frame.text = "Great Job! 🎉"
    thank_you_para = thank_you_frame.paragraphs[0]
    thank_you_para.alignment = PP_ALIGN.CENTER
    thank_you_para.font.bold = True
    thank_you_para.font.size = Pt(54)
    thank_you_para.font.color.rgb = AWS_ORANGE

    # Lesson completed
    completed = slide.shapes.add_textbox(
        Inches(1), Inches(3.8), Inches(8), Inches(1)
    )
    completed_frame = completed.text_frame
    completed_frame.text = f"You've completed: {title}"
    completed_para = completed_frame.paragraphs[0]
    completed_para.alignment = PP_ALIGN.CENTER
    completed_para.font.size = Pt(24)
    completed_para.font.color.rgb = RGBColor(255, 255, 255)

def convert_lesson_to_ppt(lesson_file, output_dir):
    """Convert a lesson JSON file to PowerPoint presentation"""

    # Load lesson data
    with open(lesson_file, 'r', encoding='utf-8') as f:
        lesson = json.load(f)

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    create_title_slide(prs, lesson['title'], lesson['icon'])

    # Process each section
    for idx, section in enumerate(lesson['sections'], 1):
        # Section header
        create_section_header(prs, idx, section['title'])

        # Section content based on type
        if section['type'] == 'text':
            create_text_slide(prs, idx, section['title'], section['content'])

        elif section['type'] == 'list':
            create_list_slide(
                prs, idx, section['title'],
                section['items'],
                section.get('content')
            )

        elif section['type'] == 'comparison':
            create_comparison_slide(
                prs, idx, section['title'],
                section['comparison']['traditional'],
                section['comparison']['cloud'],
                section.get('content')
            )

        elif section['type'] == 'detail':
            # Handle detail sections with various subsections
            content_parts = []
            if section.get('content'):
                content_parts.append(section['content'])

            # Combine all detail information
            all_items = []
            if section.get('examples'):
                all_items.extend([f"📌 {ex}" for ex in section['examples']])
            if section.get('useCases'):
                all_items.extend([f"🎯 {uc}" for uc in section['useCases']])
            if section.get('keyPoints'):
                all_items.extend([f"🔑 {kp}" for kp in section['keyPoints']])
            if section.get('bestFor'):
                all_items.extend([f"⭐ {bf}" for bf in section['bestFor']])

            if all_items:
                create_list_slide(
                    prs, idx, section['title'],
                    all_items,
                    section.get('content')
                )
            else:
                create_text_slide(prs, idx, section['title'], section['content'])

    # Advantages/Takeaways slide
    if lesson.get('advantages'):
        create_advantages_slide(prs, lesson['advantages'])

    # End slide
    create_end_slide(prs, lesson['title'])

    # Save presentation
    output_file = output_dir / f"{lesson['serviceId']}_teaching.pptx"
    prs.save(str(output_file))
    print(f"✅ Created: {output_file.name}")
    return output_file

def main():
    """Main function to convert all lessons"""

    # Setup paths
    base_dir = Path(__file__).parent
    lessons_dir = base_dir / 'data' / 'lessons'
    output_dir = base_dir / 'presentations'
    output_dir.mkdir(exist_ok=True)

    print("🎨 AWS Lesson to PowerPoint Converter")
    print("=" * 50)

    # Find all lesson JSON files
    lesson_files = list(lessons_dir.glob('**/*.json'))

    if not lesson_files:
        print("❌ No lesson files found!")
        return

    print(f"📚 Found {len(lesson_files)} lessons\n")

    # Convert each lesson
    created_files = []
    for lesson_file in lesson_files:
        try:
            output_file = convert_lesson_to_ppt(lesson_file, output_dir)
            created_files.append(output_file)
        except Exception as e:
            print(f"❌ Error converting {lesson_file.name}: {e}")

    print("\n" + "=" * 50)
    print(f"🎉 Successfully created {len(created_files)} presentations!")
    print(f"📁 Location: {output_dir}")
    print("\n📝 Presentations are ready for teaching!")
    print("   - Clean, professional design")
    print("   - AWS brand colors")
    print("   - No quizzes - just teaching content")
    print("   - Easy to present and explain")

if __name__ == '__main__':
    main()
